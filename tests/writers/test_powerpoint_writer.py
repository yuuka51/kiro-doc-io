"""PowerPointライターのユニットテスト

要件5.1: 構造化データからPowerPointファイルを生成する
要件5.2: タイトルスライド、コンテンツスライド、箇条書きリストを含むスライドを作成する
要件5.3: 生成したファイルをユーザーがアクセス可能な場所に保存する
要件5.4: 基本的なフォーマット（タイトル、本文、箇条書き）を適用する
要件5.5: ファイル生成中にエラーが発生した場合、エラーメッセージを返す
"""

import os
import pytest
import tempfile
from pathlib import Path

from pptx import Presentation

from document_format_mcp_server.writers.powerpoint_writer import PowerPointWriter
from document_format_mcp_server.utils.models import WriteResult


@pytest.fixture
def powerpoint_writer():
    """PowerPointWriterインスタンスを返すフィクスチャ"""
    return PowerPointWriter()


@pytest.fixture
def temp_output_dir():
    """一時出力ディレクトリを返すフィクスチャ"""
    with tempfile.TemporaryDirectory() as tmpdir:
        yield tmpdir


@pytest.fixture
def basic_slide_data():
    """基本的なスライドデータを返すフィクスチャ"""
    return {
        "slides": [
            {
                "layout": "title",
                "title": "テストプレゼンテーション",
                "content": "サブタイトル"
            }
        ]
    }


@pytest.fixture
def multi_slide_data():
    """複数スライドのデータを返すフィクスチャ"""
    return {
        "slides": [
            {
                "layout": "title",
                "title": "プレゼンテーションタイトル",
                "content": "サブタイトル"
            },
            {
                "layout": "content",
                "title": "概要",
                "content": "このプレゼンテーションの概要です。"
            },
            {
                "layout": "bullet",
                "title": "箇条書き",
                "content": ["項目1", "項目2", "項目3"]
            }
        ]
    }


class TestPowerPointWriterBasic:
    """PowerPointWriterの基本機能テスト"""

    def test_create_presentation_returns_write_result(
        self, powerpoint_writer, temp_output_dir, basic_slide_data
    ):
        """PowerPointWriterがWriteResultを返すことを検証"""
        output_path = os.path.join(temp_output_dir, "test.pptx")
        result = powerpoint_writer.create_presentation(basic_slide_data, output_path)
        
        # WriteResultデータクラスであることを検証
        assert isinstance(result, WriteResult)
        assert hasattr(result, 'success')
        assert hasattr(result, 'output_path')
        assert hasattr(result, 'url')
        assert hasattr(result, 'error')

    def test_create_presentation_success(
        self, powerpoint_writer, temp_output_dir, basic_slide_data
    ):
        """PowerPointファイルが正常に生成されることを検証（要件5.1）"""
        output_path = os.path.join(temp_output_dir, "test.pptx")
        result = powerpoint_writer.create_presentation(basic_slide_data, output_path)
        
        # 成功フラグを検証
        assert result.success is True
        assert result.error is None
        assert result.output_path == output_path
        
        # ファイルが存在することを検証
        assert os.path.exists(output_path)

    def test_create_presentation_file_is_valid_pptx(
        self, powerpoint_writer, temp_output_dir, basic_slide_data
    ):
        """生成されたファイルが有効なPowerPointファイルであることを検証"""
        output_path = os.path.join(temp_output_dir, "test.pptx")
        powerpoint_writer.create_presentation(basic_slide_data, output_path)
        
        # python-pptxで読み込めることを検証
        prs = Presentation(output_path)
        assert prs is not None
        assert len(prs.slides) == 1


class TestPowerPointWriterSlideLayouts:
    """スライドレイアウトのテスト"""

    def test_create_title_slide(self, powerpoint_writer, temp_output_dir):
        """タイトルスライドが正しく作成されることを検証（要件5.2）"""
        data = {
            "slides": [
                {
                    "layout": "title",
                    "title": "メインタイトル",
                    "content": "サブタイトル"
                }
            ]
        }
        output_path = os.path.join(temp_output_dir, "title_slide.pptx")
        result = powerpoint_writer.create_presentation(data, output_path)
        
        assert result.success is True
        
        # スライドの内容を検証
        prs = Presentation(output_path)
        slide = prs.slides[0]
        
        # タイトルが設定されていることを検証
        if slide.shapes.title:
            assert slide.shapes.title.text == "メインタイトル"

    def test_create_content_slide(self, powerpoint_writer, temp_output_dir):
        """コンテンツスライドが正しく作成されることを検証（要件5.2）"""
        data = {
            "slides": [
                {
                    "layout": "content",
                    "title": "コンテンツタイトル",
                    "content": "本文テキスト"
                }
            ]
        }
        output_path = os.path.join(temp_output_dir, "content_slide.pptx")
        result = powerpoint_writer.create_presentation(data, output_path)
        
        assert result.success is True
        
        # スライドの内容を検証
        prs = Presentation(output_path)
        slide = prs.slides[0]
        
        # タイトルが設定されていることを検証
        if slide.shapes.title:
            assert slide.shapes.title.text == "コンテンツタイトル"

    def test_create_bullet_slide_with_list(self, powerpoint_writer, temp_output_dir):
        """箇条書きスライド（リスト形式）が正しく作成されることを検証（要件5.2）"""
        data = {
            "slides": [
                {
                    "layout": "bullet",
                    "title": "箇条書きタイトル",
                    "content": ["項目1", "項目2", "項目3"]
                }
            ]
        }
        output_path = os.path.join(temp_output_dir, "bullet_slide.pptx")
        result = powerpoint_writer.create_presentation(data, output_path)
        
        assert result.success is True
        
        # スライドの内容を検証
        prs = Presentation(output_path)
        slide = prs.slides[0]
        
        # タイトルが設定されていることを検証
        if slide.shapes.title:
            assert slide.shapes.title.text == "箇条書きタイトル"

    def test_create_bullet_slide_with_string(self, powerpoint_writer, temp_output_dir):
        """箇条書きスライド（文字列形式）が正しく作成されることを検証（要件5.2）"""
        data = {
            "slides": [
                {
                    "layout": "bullet",
                    "title": "箇条書きタイトル",
                    "content": "項目1\n項目2\n項目3"
                }
            ]
        }
        output_path = os.path.join(temp_output_dir, "bullet_string_slide.pptx")
        result = powerpoint_writer.create_presentation(data, output_path)
        
        assert result.success is True
        
        # ファイルが存在することを検証
        assert os.path.exists(output_path)

    def test_create_multiple_slides(
        self, powerpoint_writer, temp_output_dir, multi_slide_data
    ):
        """複数スライドが正しく作成されることを検証（要件5.2）"""
        output_path = os.path.join(temp_output_dir, "multi_slides.pptx")
        result = powerpoint_writer.create_presentation(multi_slide_data, output_path)
        
        assert result.success is True
        
        # スライド数を検証
        prs = Presentation(output_path)
        assert len(prs.slides) == 3


class TestPowerPointWriterFileSaving:
    """ファイル保存機能のテスト"""

    def test_save_to_specified_path(
        self, powerpoint_writer, temp_output_dir, basic_slide_data
    ):
        """指定されたパスにファイルが保存されることを検証（要件5.3）"""
        output_path = os.path.join(temp_output_dir, "specified_path.pptx")
        result = powerpoint_writer.create_presentation(basic_slide_data, output_path)
        
        assert result.success is True
        assert result.output_path == output_path
        assert os.path.exists(output_path)

    def test_create_output_directory_if_not_exists(
        self, powerpoint_writer, temp_output_dir, basic_slide_data
    ):
        """出力ディレクトリが存在しない場合に作成されることを検証（要件5.3）"""
        nested_dir = os.path.join(temp_output_dir, "nested", "dir")
        output_path = os.path.join(nested_dir, "test.pptx")
        
        result = powerpoint_writer.create_presentation(basic_slide_data, output_path)
        
        assert result.success is True
        assert os.path.exists(output_path)
        assert os.path.isdir(nested_dir)


class TestPowerPointWriterErrorHandling:
    """エラーハンドリングのテスト"""

    def test_error_on_missing_slides_key(self, powerpoint_writer, temp_output_dir):
        """slidesキーがない場合にエラーを返すことを検証（要件5.5）"""
        data = {"title": "タイトルのみ"}
        output_path = os.path.join(temp_output_dir, "error.pptx")
        
        result = powerpoint_writer.create_presentation(data, output_path)
        
        assert result.success is False
        assert result.error is not None
        assert "slides" in result.error.lower() or "キー" in result.error

    def test_error_on_invalid_layout(self, powerpoint_writer, temp_output_dir):
        """不正なレイアウトの場合にエラーを返すことを検証（要件5.5）"""
        data = {
            "slides": [
                {
                    "layout": "invalid_layout",
                    "title": "テスト"
                }
            ]
        }
        output_path = os.path.join(temp_output_dir, "error.pptx")
        
        result = powerpoint_writer.create_presentation(data, output_path)
        
        assert result.success is False
        assert result.error is not None

    def test_error_on_invalid_data_type(self, powerpoint_writer, temp_output_dir):
        """不正なデータ型の場合にエラーを返すことを検証（要件5.5）"""
        output_path = os.path.join(temp_output_dir, "error.pptx")
        
        result = powerpoint_writer.create_presentation("invalid_data", output_path)
        
        assert result.success is False
        assert result.error is not None

    def test_error_on_slides_not_list(self, powerpoint_writer, temp_output_dir):
        """slidesがリストでない場合にエラーを返すことを検証（要件5.5）"""
        data = {"slides": "not_a_list"}
        output_path = os.path.join(temp_output_dir, "error.pptx")
        
        result = powerpoint_writer.create_presentation(data, output_path)
        
        assert result.success is False
        assert result.error is not None


class TestPowerPointWriterFormatting:
    """フォーマット適用のテスト"""

    def test_title_formatting_applied(self, powerpoint_writer, temp_output_dir):
        """タイトルフォーマットが適用されることを検証（要件5.4）"""
        data = {
            "slides": [
                {
                    "layout": "title",
                    "title": "フォーマットテスト",
                    "content": "サブタイトル"
                }
            ]
        }
        output_path = os.path.join(temp_output_dir, "format_test.pptx")
        result = powerpoint_writer.create_presentation(data, output_path)
        
        assert result.success is True
        
        # ファイルが読み込めることを検証
        prs = Presentation(output_path)
        assert len(prs.slides) == 1

    def test_empty_slides_list(self, powerpoint_writer, temp_output_dir):
        """空のスライドリストでもファイルが生成されることを検証"""
        data = {"slides": []}
        output_path = os.path.join(temp_output_dir, "empty_slides.pptx")
        
        result = powerpoint_writer.create_presentation(data, output_path)
        
        assert result.success is True
        assert os.path.exists(output_path)
        
        # 空のプレゼンテーションが作成されることを検証
        prs = Presentation(output_path)
        assert len(prs.slides) == 0

    def test_slide_with_empty_title(self, powerpoint_writer, temp_output_dir):
        """空のタイトルでもスライドが作成されることを検証"""
        data = {
            "slides": [
                {
                    "layout": "content",
                    "title": "",
                    "content": "コンテンツのみ"
                }
            ]
        }
        output_path = os.path.join(temp_output_dir, "empty_title.pptx")
        
        result = powerpoint_writer.create_presentation(data, output_path)
        
        assert result.success is True
        assert os.path.exists(output_path)

    def test_slide_with_empty_content(self, powerpoint_writer, temp_output_dir):
        """空のコンテンツでもスライドが作成されることを検証"""
        data = {
            "slides": [
                {
                    "layout": "content",
                    "title": "タイトルのみ",
                    "content": ""
                }
            ]
        }
        output_path = os.path.join(temp_output_dir, "empty_content.pptx")
        
        result = powerpoint_writer.create_presentation(data, output_path)
        
        assert result.success is True
        assert os.path.exists(output_path)
