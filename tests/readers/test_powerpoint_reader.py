"""PowerPointリーダーのユニットテスト（tests/readers/配下）

要件1.1: PowerPointファイルからテキスト内容、スライド構造、画像の説明を抽出する
要件1.2: 抽出したコンテンツを構造化されたテキスト形式でKiroに提供する
要件1.3: PowerPointファイルが破損しているまたは読み取り不可能な場合、エラーメッセージを返す
要件1.4: 各スライドのタイトル、本文、ノート、表データを識別して抽出する

このファイルはtests/test_powerpoint_reader.pyから全てのテストをインポートし、
追加のテストケースを提供します。
"""

import pytest
import os
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from document_format_mcp_server.readers.powerpoint_reader import PowerPointReader
from document_format_mcp_server.utils.models import ReadResult, DocumentContent

# tests/test_powerpoint_reader.pyから全てのテストをインポート
from tests.test_powerpoint_reader import *


# ===== 追加のテストケース =====

@pytest.fixture
def temp_pptx_with_notes():
    """ノート付きのPowerPointファイルを作成するフィクスチャ"""
    prs = Presentation()
    
    # タイトルスライドを追加
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "テストプレゼンテーション"
    subtitle.text = "サブタイトル"
    
    # ノートを追加
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "これはスライド1のノートです。"
    
    # コンテンツスライドを追加
    content_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_layout)
    slide2.shapes.title.text = "コンテンツスライド"
    
    # ノートを追加
    notes_slide2 = slide2.notes_slide
    notes_slide2.notes_text_frame.text = "これはスライド2のノートです。"
    
    # 一時ファイルに保存
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        temp_path = f.name
    
    prs.save(temp_path)
    
    yield temp_path
    
    # クリーンアップ
    if os.path.exists(temp_path):
        os.unlink(temp_path)


@pytest.fixture
def temp_pptx_with_table():
    """表付きのPowerPointファイルを作成するフィクスチャ"""
    prs = Presentation()
    
    # 空白スライドを追加
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 表を追加
    rows, cols = 3, 3
    left = Inches(1)
    top = Inches(1)
    width = Inches(6)
    height = Inches(2)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # ヘッダー行
    table.cell(0, 0).text = "項目"
    table.cell(0, 1).text = "値"
    table.cell(0, 2).text = "備考"
    
    # データ行
    table.cell(1, 0).text = "項目1"
    table.cell(1, 1).text = "100"
    table.cell(1, 2).text = "テスト"
    
    table.cell(2, 0).text = "項目2"
    table.cell(2, 1).text = "200"
    table.cell(2, 2).text = "サンプル"
    
    # 一時ファイルに保存
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        temp_path = f.name
    
    prs.save(temp_path)
    
    yield temp_path
    
    # クリーンアップ
    if os.path.exists(temp_path):
        os.unlink(temp_path)


class TestPowerPointReaderNotes:
    """ノート抽出のテスト"""

    def test_extracts_slide_notes(self, temp_pptx_with_notes):
        """スライドノートが正しく抽出されることを検証（要件1.4）"""
        reader = PowerPointReader()
        result = reader.read_file(temp_pptx_with_notes)
        
        assert result.success is True
        slides = result.content.content["slides"]
        
        # 少なくとも1つのスライドにノートがあることを検証
        has_notes = False
        for slide in slides:
            if slide["notes"] and len(slide["notes"]) > 0:
                has_notes = True
                break
        
        assert has_notes, "少なくとも1つのスライドにノートが必要です"

    def test_notes_content_is_correct(self, temp_pptx_with_notes):
        """ノートの内容が正しいことを検証（要件1.4）"""
        reader = PowerPointReader()
        result = reader.read_file(temp_pptx_with_notes)
        
        assert result.success is True
        slides = result.content.content["slides"]
        
        # 最初のスライドのノートを検証
        first_slide = slides[0]
        assert "ノート" in first_slide["notes"] or "notes" in first_slide["notes"].lower()


class TestPowerPointReaderTables:
    """表抽出のテスト"""

    def test_extracts_table_data(self, temp_pptx_with_table):
        """表データが正しく抽出されることを検証（要件1.4）"""
        reader = PowerPointReader()
        result = reader.read_file(temp_pptx_with_table)
        
        assert result.success is True
        slides = result.content.content["slides"]
        
        # 表を含むスライドを探す
        table_found = False
        for slide in slides:
            if slide["tables"]:
                table_found = True
                table = slide["tables"][0]
                
                # 表の構造を検証
                assert "rows" in table
                assert "columns" in table
                assert "data" in table
                assert table["rows"] == 3
                assert table["columns"] == 3
                break
        
        assert table_found, "表を含むスライドが見つかりません"

    def test_table_data_content(self, temp_pptx_with_table):
        """表のデータ内容が正しいことを検証（要件1.4）"""
        reader = PowerPointReader()
        result = reader.read_file(temp_pptx_with_table)
        
        assert result.success is True
        slides = result.content.content["slides"]
        
        # 表を含むスライドを探す
        for slide in slides:
            if slide["tables"]:
                table = slide["tables"][0]
                data = table["data"]
                
                # ヘッダー行を検証
                assert "項目" in data[0]
                assert "値" in data[0]
                assert "備考" in data[0]
                
                # データ行を検証
                assert "項目1" in data[1]
                assert "100" in data[1]
                break


class TestPowerPointReaderStructuredOutput:
    """構造化出力のテスト"""

    def test_output_is_structured_format(self, temp_pptx_with_notes):
        """出力が構造化された形式であることを検証（要件1.2）"""
        reader = PowerPointReader()
        result = reader.read_file(temp_pptx_with_notes)
        
        assert result.success is True
        
        # DocumentContentの構造を検証
        assert isinstance(result.content, DocumentContent)
        assert result.content.format_type == "pptx"
        
        # メタデータの構造を検証
        metadata = result.content.metadata
        assert isinstance(metadata, dict)
        assert "slide_count" in metadata
        assert "total_slides" in metadata
        assert "file_size_mb" in metadata
        
        # コンテンツの構造を検証
        content = result.content.content
        assert isinstance(content, dict)
        assert "slides" in content
        assert isinstance(content["slides"], list)

    def test_each_slide_has_required_fields(self, temp_pptx_with_notes):
        """各スライドに必要なフィールドがあることを検証（要件1.4）"""
        reader = PowerPointReader()
        result = reader.read_file(temp_pptx_with_notes)
        
        assert result.success is True
        slides = result.content.content["slides"]
        
        required_fields = ["slide_number", "title", "content", "notes", "tables"]
        
        for slide in slides:
            for field in required_fields:
                assert field in slide, f"スライドに'{field}'フィールドがありません"


class TestPowerPointReaderErrorHandling:
    """エラーハンドリングの追加テスト"""

    def test_invalid_extension_file(self):
        """不正な拡張子のファイルを読み取った場合のエラーハンドリングを検証"""
        reader = PowerPointReader()
        
        # 一時的なテキストファイルを作成
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
            f.write("これはテキストファイルです")
            temp_path = f.name
        
        try:
            result = reader.read_file(temp_path)
            
            # 失敗フラグを検証
            assert result.success is False
            assert result.error is not None
        finally:
            if os.path.exists(temp_path):
                os.unlink(temp_path)

    def test_empty_pptx_file(self):
        """空のPowerPointファイルを読み取った場合の処理を検証"""
        reader = PowerPointReader()
        
        # 空のPowerPointファイルを作成
        prs = Presentation()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            temp_path = f.name
        
        prs.save(temp_path)
        
        try:
            result = reader.read_file(temp_path)
            
            # 成功するが、スライドが空であることを検証
            assert result.success is True
            assert result.content.content["slides"] == []
            assert result.content.metadata["slide_count"] == 0
        finally:
            if os.path.exists(temp_path):
                os.unlink(temp_path)


class TestPowerPointReaderConfiguration:
    """設定パラメータのテスト"""

    def test_default_configuration(self):
        """デフォルト設定でPowerPointReaderを初期化できることを検証"""
        reader = PowerPointReader()
        
        assert reader.max_slides == 500
        assert reader.max_file_size_mb == 100

    def test_custom_max_slides(self):
        """カスタムmax_slidesパラメータでPowerPointReaderを初期化できることを検証"""
        reader = PowerPointReader(max_slides=10)
        
        assert reader.max_slides == 10

    def test_custom_max_file_size(self):
        """カスタムmax_file_size_mbパラメータでPowerPointReaderを初期化できることを検証"""
        reader = PowerPointReader(max_file_size_mb=50)
        
        assert reader.max_file_size_mb == 50

    def test_custom_both_parameters(self):
        """両方のカスタムパラメータでPowerPointReaderを初期化できることを検証"""
        reader = PowerPointReader(max_slides=20, max_file_size_mb=25)
        
        assert reader.max_slides == 20
        assert reader.max_file_size_mb == 25
