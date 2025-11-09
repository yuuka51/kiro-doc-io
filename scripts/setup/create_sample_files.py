"""テスト用のサンプルファイルを生成するスクリプト"""

import sys
from pathlib import Path

# srcディレクトリをパスに追加
sys.path.insert(0, str(Path(__file__).parent / "src"))

from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches as DocxInches
from openpyxl import Workbook


def create_sample_powerpoint(output_path: str):
    """サンプルPowerPointファイルを作成"""
    print(f"PowerPointファイルを作成中: {output_path}")
    
    prs = Presentation()
    
    # スライド1: タイトルスライド
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "サンプルプレゼンテーション"
    subtitle.text = "Document Format MCP Server テスト用"
    
    # スライド2: タイトルとコンテンツ
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    title.text = "機能紹介"
    
    tf = content.text_frame
    tf.text = "主な機能:"
    
    p = tf.add_paragraph()
    p.text = "PowerPointファイルの読み取り"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Wordファイルの読み取り"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Excelファイルの読み取り"
    p.level = 1
    
    # スライド3: 表を含むスライド
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide3.shapes.title
    title.text = "データ表"
    
    # 表を追加
    rows, cols = 4, 3
    left = Inches(2.0)
    top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(2.0)
    
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table
    
    # ヘッダー行
    table.cell(0, 0).text = "項目"
    table.cell(0, 1).text = "値"
    table.cell(0, 2).text = "備考"
    
    # データ行
    table.cell(1, 0).text = "読み取り"
    table.cell(1, 1).text = "対応"
    table.cell(1, 2).text = "完了"
    
    table.cell(2, 0).text = "書き込み"
    table.cell(2, 1).text = "対応予定"
    table.cell(2, 2).text = "開発中"
    
    table.cell(3, 0).text = "変換"
    table.cell(3, 1).text = "対応予定"
    table.cell(3, 2).text = "計画中"
    
    prs.save(output_path)
    print(f"✅ 作成完了: {output_path}")


def create_sample_word(output_path: str):
    """サンプルWordファイルを作成"""
    print(f"Wordファイルを作成中: {output_path}")
    
    doc = Document()
    
    # タイトル
    doc.add_heading('サンプルドキュメント', 0)
    
    # 段落
    doc.add_paragraph('これはDocument Format MCP Serverのテスト用サンプルファイルです。')
    
    # 見出し1
    doc.add_heading('概要', level=1)
    doc.add_paragraph(
        'このMCPサーバは、Microsoft OfficeおよびGoogle Workspaceのドキュメント形式を'
        '読み書きするための機能を提供します。'
    )
    
    # 見出し2
    doc.add_heading('対応フォーマット', level=2)
    
    # 箇条書き
    doc.add_paragraph('PowerPoint (.pptx)', style='List Bullet')
    doc.add_paragraph('Word (.docx)', style='List Bullet')
    doc.add_paragraph('Excel (.xlsx)', style='List Bullet')
    doc.add_paragraph('Google スプレッドシート', style='List Bullet')
    doc.add_paragraph('Google ドキュメント', style='List Bullet')
    doc.add_paragraph('Google スライド', style='List Bullet')
    
    # 見出し2
    doc.add_heading('機能', level=2)
    
    # 番号付きリスト
    doc.add_paragraph('ファイルの読み取り', style='List Number')
    doc.add_paragraph('ファイルの書き込み', style='List Number')
    doc.add_paragraph('フォーマット変換', style='List Number')
    
    # 表
    doc.add_heading('ステータス表', level=2)
    table = doc.add_table(rows=4, cols=3)
    table.style = 'Light Grid Accent 1'
    
    # ヘッダー行
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = '機能'
    hdr_cells[1].text = 'ステータス'
    hdr_cells[2].text = '備考'
    
    # データ行
    row1_cells = table.rows[1].cells
    row1_cells[0].text = 'ローカルファイル読み取り'
    row1_cells[1].text = '完了'
    row1_cells[2].text = 'PowerPoint, Word, Excel'
    
    row2_cells = table.rows[2].cells
    row2_cells[0].text = 'Google Workspace読み取り'
    row2_cells[1].text = '完了'
    row2_cells[2].text = 'スプレッドシート, ドキュメント, スライド'
    
    row3_cells = table.rows[3].cells
    row3_cells[0].text = 'ファイル書き込み'
    row3_cells[1].text = '開発中'
    row3_cells[2].text = 'タスク5-7で実装予定'
    
    doc.save(output_path)
    print(f"✅ 作成完了: {output_path}")


def create_sample_excel(output_path: str):
    """サンプルExcelファイルを作成"""
    print(f"Excelファイルを作成中: {output_path}")
    
    wb = Workbook()
    
    # シート1: データ
    ws1 = wb.active
    ws1.title = "データ"
    
    # ヘッダー
    ws1['A1'] = "ID"
    ws1['B1'] = "名前"
    ws1['C1'] = "カテゴリ"
    ws1['D1'] = "値"
    
    # データ
    data = [
        [1, "PowerPoint", "読み取り", "完了"],
        [2, "Word", "読み取り", "完了"],
        [3, "Excel", "読み取り", "完了"],
        [4, "Google Sheets", "読み取り", "完了"],
        [5, "Google Docs", "読み取り", "完了"],
        [6, "Google Slides", "読み取り", "完了"],
    ]
    
    for row_idx, row_data in enumerate(data, start=2):
        for col_idx, value in enumerate(row_data, start=1):
            ws1.cell(row=row_idx, column=col_idx, value=value)
    
    # シート2: 統計
    ws2 = wb.create_sheet("統計")
    
    ws2['A1'] = "項目"
    ws2['B1'] = "数"
    
    ws2['A2'] = "対応フォーマット"
    ws2['B2'] = 6
    
    ws2['A3'] = "完了タスク"
    ws2['B3'] = 4
    
    ws2['A4'] = "残りタスク"
    ws2['B4'] = 6
    
    # シート3: 計算
    ws3 = wb.create_sheet("計算")
    
    ws3['A1'] = "項目"
    ws3['B1'] = "値"
    ws3['C1'] = "計算"
    
    ws3['A2'] = "数値1"
    ws3['B2'] = 100
    
    ws3['A3'] = "数値2"
    ws3['B3'] = 50
    
    ws3['A4'] = "合計"
    ws3['C4'] = "=B2+B3"
    
    ws3['A5'] = "平均"
    ws3['C5'] = "=(B2+B3)/2"
    
    wb.save(output_path)
    print(f"✅ 作成完了: {output_path}")


def main():
    """メイン関数"""
    print("\n" + "="*60)
    print("サンプルファイル生成スクリプト")
    print("="*60 + "\n")
    
    # test_filesディレクトリを作成
    test_dir = Path("test_files")
    test_dir.mkdir(exist_ok=True)
    print(f"出力ディレクトリ: {test_dir}\n")
    
    # サンプルファイルを生成
    try:
        create_sample_powerpoint(str(test_dir / "sample.pptx"))
        create_sample_word(str(test_dir / "sample.docx"))
        create_sample_excel(str(test_dir / "sample.xlsx"))
        
        print("\n" + "="*60)
        print("✅ すべてのサンプルファイルの生成が完了しました")
        print("="*60)
        print("\n次のコマンドでテストを実行できます:")
        print("  python test_readers.py")
        print("="*60 + "\n")
        
    except Exception as e:
        print(f"\n❌ エラーが発生しました: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
