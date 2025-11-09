"""PowerPoint (.pptx) file reader."""

import os
import time
from typing import Any

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from ..utils.errors import FileNotFoundError, CorruptedFileError
from ..utils.logging_config import get_logger


# ロガーの取得
logger = get_logger(__name__)


class PowerPointReader:
    """PowerPoint (.pptx)ファイルを読み取るクラス。"""

    def __init__(self, max_slides: int = 500, max_file_size_mb: int = 100):
        """
        PowerPointリーダーを初期化する。

        Args:
            max_slides: 処理する最大スライド数（デフォルト: 500）
            max_file_size_mb: 処理する最大ファイルサイズ（MB）（デフォルト: 100）
        """
        self.max_slides = max_slides
        self.max_file_size_mb = max_file_size_mb

    def read_file(self, file_path: str) -> dict[str, Any]:
        """
        PowerPointファイルからコンテンツを抽出する。

        Args:
            file_path: 読み取るPowerPointファイルのパス

        Returns:
            抽出されたコンテンツを含む辞書:
            {
                "slides": [
                    {
                        "slide_number": int,
                        "title": str,
                        "content": str,
                        "notes": str,
                        "tables": [...]
                    }
                ]
            }

        Raises:
            FileNotFoundError: ファイルが存在しない場合
            CorruptedFileError: ファイルが破損している場合
        """
        logger.info(f"PowerPointファイルの読み込みを開始: {file_path}")
        start_time = time.time()
        
        # ファイルの存在確認
        if not os.path.exists(file_path):
            logger.error(f"ファイルが見つかりません: {file_path}")
            raise FileNotFoundError(
                f"指定されたファイルが見つかりません: {file_path}",
                details={"file_path": file_path}
            )
        
        # ファイルサイズの検証
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        if file_size_mb > self.max_file_size_mb:
            logger.error(
                f"ファイルサイズが制限を超えています: {file_size_mb:.2f}MB > {self.max_file_size_mb}MB"
            )
            raise CorruptedFileError(
                f"ファイルサイズが制限を超えています: {file_size_mb:.2f}MB（最大: {self.max_file_size_mb}MB）",
                details={
                    "file_path": file_path,
                    "file_size_mb": file_size_mb,
                    "max_file_size_mb": self.max_file_size_mb
                }
            )

        try:
            # PowerPointファイルを開く
            prs = Presentation(file_path)
            
            # スライド数の検証
            slide_count = len(prs.slides)
            if slide_count > self.max_slides:
                logger.warning(
                    f"スライド数が制限を超えています: {slide_count} > {self.max_slides}。"
                    f"最初の{self.max_slides}スライドのみを処理します。"
                )
            
            slides_data = []
            
            # 各スライドを処理（制限まで）
            slides_to_process = min(slide_count, self.max_slides)
            for idx, slide in enumerate(list(prs.slides)[:slides_to_process], start=1):
                slide_data = {
                    "slide_number": idx,
                    "title": self._extract_title(slide),
                    "content": self._extract_content(slide),
                    "notes": self._extract_notes(slide),
                    "tables": self._extract_tables(slide)
                }
                slides_data.append(slide_data)
            
            result = {"slides": slides_data}
            
            # スライド数制限の警告を追加
            if slide_count > self.max_slides:
                result["warning"] = f"ファイルには{slide_count}個のスライドがありますが、最初の{self.max_slides}スライドのみを処理しました。"
            
            # 処理時間を計算
            elapsed_time = time.time() - start_time
            logger.info(
                f"PowerPointファイルの読み込みが完了: {file_path} "
                f"(スライド数: {len(slides_data)}, 処理時間: {elapsed_time:.2f}秒)"
            )
            logger.debug(f"抽出されたデータの概要: スライド数={len(slides_data)}")
            
            return result
        
        except PackageNotFoundError as e:
            logger.error(
                f"PowerPointファイルが破損しています: {file_path}",
                exc_info=True
            )
            raise CorruptedFileError(
                f"PowerPointファイルが破損しているか、読み取り不可能です: {file_path}",
                details={"file_path": file_path, "error": str(e)}
            )
        except Exception as e:
            # その他の予期しないエラー
            logger.error(
                f"PowerPointファイルの読み取り中にエラーが発生: {file_path}",
                exc_info=True
            )
            raise CorruptedFileError(
                f"PowerPointファイルの読み取り中にエラーが発生しました: {file_path}",
                details={"file_path": file_path, "error": str(e)}
            )

    def _extract_title(self, slide) -> str:
        """スライドからタイトルを抽出する。"""
        if slide.shapes.title:
            return slide.shapes.title.text
        return ""

    def _extract_content(self, slide) -> str:
        """スライドから本文コンテンツを抽出する。"""
        content_parts = []
        
        for shape in slide.shapes:
            # タイトル以外のテキストフレームを処理
            if hasattr(shape, "text") and shape != slide.shapes.title:
                text = shape.text.strip()
                if text:
                    content_parts.append(text)
        
        return "\n\n".join(content_parts)

    def _extract_notes(self, slide) -> str:
        """スライドからノートを抽出する。"""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                return notes_slide.notes_text_frame.text.strip()
        return ""

    def _extract_tables(self, slide) -> list[dict[str, Any]]:
        """スライドから表データを抽出する。"""
        tables = []
        
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_data = {
                    "rows": len(table.rows),
                    "columns": len(table.columns),
                    "data": []
                }
                
                # 各行のデータを抽出
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data["data"].append(row_data)
                
                tables.append(table_data)
        
        return tables
