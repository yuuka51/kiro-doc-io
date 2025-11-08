"""PowerPoint (.pptx) file writer."""

import os
import time
from typing import Any

from pptx import Presentation

from ..utils.errors import ValidationError
from ..utils.logging_config import get_logger


# ロガーの取得
logger = get_logger(__name__)


class PowerPointWriter:
    """PowerPoint (.pptx)ファイルを生成するクラス。"""

    def create_presentation(self, data: dict[str, Any], output_path: str) -> str:
        """
        構造化データからPowerPointファイルを生成する。

        Args:
            data: プレゼンテーションデータ:
                {
                    "title": str (オプション),
                    "slides": [
                        {
                            "layout": "title" | "content" | "bullet",
                            "title": str,
                            "content": str | list
                        }
                    ]
                }
            output_path: 出力ファイルのパス

        Returns:
            作成されたファイルのパス

        Raises:
            ValidationError: 入力データが不正な場合
            Exception: ファイル生成中にエラーが発生した場合
        """
        logger.info(f"PowerPointファイルの生成を開始: {output_path}")
        start_time = time.time()

        try:
            # 入力データの検証
            self._validate_data(data)

            # 新しいプレゼンテーションを作成
            prs = Presentation()

            # スライドを作成
            slides_data = data.get("slides", [])
            for slide_data in slides_data:
                self._create_slide(prs, slide_data)

            # 出力ディレクトリが存在しない場合は作成
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                try:
                    os.makedirs(output_dir, exist_ok=True)
                    logger.debug(f"出力ディレクトリを作成: {output_dir}")
                except OSError as e:
                    logger.error(f"出力ディレクトリの作成に失敗: {output_dir}", exc_info=True)
                    raise Exception(
                        f"出力ディレクトリの作成に失敗しました: {output_dir}"
                    ) from e

            # ファイルを保存
            try:
                prs.save(output_path)
            except PermissionError as e:
                logger.error(f"ファイルの保存に失敗（権限エラー）: {output_path}", exc_info=True)
                raise Exception(
                    f"ファイルの保存に失敗しました（権限エラー）: {output_path}"
                ) from e
            except OSError as e:
                logger.error(f"ファイルの保存に失敗: {output_path}", exc_info=True)
                raise Exception(
                    f"ファイルの保存に失敗しました: {output_path}"
                ) from e

            # 処理時間を計算
            elapsed_time = time.time() - start_time
            logger.info(
                f"PowerPointファイルの生成が完了: {output_path} "
                f"(スライド数: {len(slides_data)}, 処理時間: {elapsed_time:.2f}秒)"
            )

            return output_path

        except ValidationError:
            # ValidationErrorはそのまま再送出
            raise
        except Exception as e:
            logger.error(
                f"PowerPointファイルの生成中にエラーが発生: {output_path}",
                exc_info=True
            )
            raise Exception(
                f"PowerPointファイルの生成中にエラーが発生しました: {str(e)}"
            ) from e

    def _validate_data(self, data: dict[str, Any]) -> None:
        """入力データを検証する。"""
        if not isinstance(data, dict):
            raise ValidationError(
                "入力データは辞書形式である必要があります",
                details={"data_type": type(data).__name__}
            )

        if "slides" not in data:
            raise ValidationError(
                "入力データに'slides'キーが必要です",
                details={"keys": list(data.keys())}
            )

        slides = data["slides"]
        if not isinstance(slides, list):
            raise ValidationError(
                "'slides'はリスト形式である必要があります",
                details={"slides_type": type(slides).__name__}
            )

        # 各スライドデータを検証
        for idx, slide_data in enumerate(slides):
            if not isinstance(slide_data, dict):
                raise ValidationError(
                    f"スライド{idx + 1}のデータは辞書形式である必要があります",
                    details={"slide_index": idx, "data_type": type(slide_data).__name__}
                )

            layout = slide_data.get("layout", "content")
            if layout not in ["title", "content", "bullet"]:
                raise ValidationError(
                    f"スライド{idx + 1}のレイアウトが不正です: {layout}",
                    details={"slide_index": idx, "layout": layout}
                )

    def _create_slide(self, prs: Presentation, slide_data: dict[str, Any]) -> None:
        """スライドを作成する。"""
        layout_type = slide_data.get("layout", "content")
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")

        if layout_type == "title":
            self._create_title_slide(prs, title, content)
        elif layout_type == "bullet":
            self._create_bullet_slide(prs, title, content)
        else:  # content
            self._create_content_slide(prs, title, content)

    def _create_title_slide(self, prs: Presentation, title: str, subtitle: str) -> None:
        """タイトルスライドを作成する。"""
        # タイトルスライドレイアウト（通常は0番目）
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # タイトルとサブタイトルを設定
        if slide.shapes.title:
            slide.shapes.title.text = title

        # サブタイトルを設定（プレースホルダーが存在する場合）
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        logger.debug(f"タイトルスライドを作成: {title}")

    def _create_content_slide(self, prs: Presentation, title: str, content: str) -> None:
        """コンテンツスライドを作成する。"""
        # タイトルとコンテンツのレイアウト（通常は1番目）
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # タイトルを設定
        if slide.shapes.title:
            slide.shapes.title.text = title

        # コンテンツを設定
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.text = str(content)

        logger.debug(f"コンテンツスライドを作成: {title}")

    def _create_bullet_slide(
        self, prs: Presentation, title: str, content: str | list
    ) -> None:
        """箇条書きスライドを作成する。"""
        # タイトルとコンテンツのレイアウト（通常は1番目）
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # タイトルを設定
        if slide.shapes.title:
            slide.shapes.title.text = title

        # 箇条書きコンテンツを設定
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # 既存のテキストをクリア

            # contentがリストの場合は各項目を箇条書きに
            if isinstance(content, list):
                for idx, item in enumerate(content):
                    if idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = str(item)
                    p.level = 0
            else:
                # 文字列の場合は改行で分割して箇条書きに
                lines = str(content).split("\n")
                for idx, line in enumerate(lines):
                    line = line.strip()
                    if not line:
                        continue
                    if idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = line
                    p.level = 0

        logger.debug(f"箇条書きスライドを作成: {title}")
